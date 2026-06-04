#!/usr/bin/env python3
"""
RTTL Complete Strategy PowerPoint - Final Version
Think! Design & Planning, LLC
"""

from pptx import Presentation
from pptx.util import Pt

TEMPLATE = r"C:\Users\Chris\Downloads\NCATSTAT_Template.pptx"
OUTPUT   = r"C:\Users\Chris\Downloads\RTTL-Business-Research-Think.pptx"

prs = Presentation(TEMPLATE)

# Delete all example slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

L_TITLE = prs.slide_layouts[0]
L_TRANS = prs.slide_layouts[2]
L1 = prs.slide_layouts[5]
L2 = prs.slide_layouts[6]
L3 = prs.slide_layouts[7]
L4 = prs.slide_layouts[8]
layouts = [L1, L2, L3, L4]
li = 0

def next_layout():
    global li
    layout = layouts[li % 4]
    li += 1
    return layout

def transition(text):
    s = prs.slides.add_slide(L_TRANS)
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 10:
            ph.text = text
    return s

def content(title, bullets):
    s = prs.slides.add_slide(next_layout())
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            break
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            for i, (t, lv) in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = t
                p.level = lv
                for r in p.runs:
                    r.font.size = Pt(13)
            break
    return s

# =====================================================================
# SLIDE 1 - Title
# =====================================================================
s = prs.slides.add_slide(L_TITLE)
for ph in s.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "RTTL Complete Strategy"
print("1: Title")

# =====================================================================
# PART 1: THE VISION
# =====================================================================
transition("Part 1: The Vision\nIntegrated Agrotherapy Campus")
print("2: Vision transition")

content("The Big Idea", [
    ("The elderly care facility is the KEYSTONE", 0),
    ("It makes the entire 66-acre estate self-sustaining", 1),
    ("The lavender farm = therapeutic infrastructure", 0),
    ("Training programs produce workers for farm AND facility", 0),
    ("The farm feeds the facility", 0),
    ("The facility funds the farm", 0),
    ("The training programs staff both", 0),
])
print("3: Big Idea")

content("Lavender & Hemp -- Elder Care Medicine", [
    ("Dementia Agitation -- Lavender aromatherapy", 0),
    ("Chronic Pain -- Lavender oil AND hemp creams for joints", 0),
    ("Insomnia -- Lavender oil improves sleep in 65+", 0),
    ("Muscle & Joint -- Hemp-infused topicals, made on-site", 0),
    ("Stress / Soul Healing -- Personal CEDAR SAUNAS", 0),
    ("Infrared heat therapy: circulation, detox, relaxation", 1),
    ("Cognitive Decline -- Sensory garden stimulation", 0),
])
print("4: Lavender & Hemp Medicine")

content("10 Dual-Purpose Therapy Areas", [
    ("Plant Processing -- Residents harvest & bundle lavender", 0),
    ("Animal Husbandry -- Egg collection, animal-assisted therapy", 0),
    ("Hemp Wellness -- Residents make hemp creams they use daily", 0),
    ("Cedar Saunas -- Personal sessions for soul healing", 0),
    ("Food Processing -- Farm-to-table cooking programs", 0),
    ("Lavender Products -- Bundling sachets, labeling soaps", 0),
    ("Conservation -- Bird watching, pollinator gardens", 0),
    ("Education / History -- AKA legacy storytelling circles", 0),
])
print("5: 10 Therapy Areas")

content("Intergenerational Programming", [
    ("Young people and elderly work TOGETHER, not separately", 0),
    ("Lavender Harvest: Youth harvest, Elders bundle & teach", 0),
    ("Hemp Cream Production: Youth process, Elders test & advise", 0),
    ("Cedar Sauna Sessions: Youth manage, Elders rest & heal", 0),
    ("Animal Care: Youth feed/clean, Elders collect eggs & bond", 0),
    ("Farm Store: Youth do e-commerce, Elders greet & tell stories", 0),
    ("Cooking: Youth prep & serve, Elders share recipes & wisdom", 0),
    ("This is how farming communities have always worked", 1),
])
print("6: Intergenerational")

content("Hemp R&D -- 5 Acres Dedicated", [
    ("Hemp field = R&D platform, not just a crop", 0),
    ("Research alternative uses until full legalization", 1),
    ("Topical Creams -- Pain relief, moisturizers, muscle rubs", 0),
    ("Produced on-site, used by residents daily", 1),
    ("Construction -- Hempcrete blocks for future ALF expansion", 0),
    ("Textiles -- Bedding & towels from estate-grown hemp", 0),
    ("Animal Bedding -- Hemp hurd for goats & chickens", 0),
    ("Bioplastics -- Sustainable packaging for farm store", 0),
])
print("7: Hemp R&D")

content("Farm-to-Table Nutrition -- Food as Medicine", [
    ("Vegetable-based nutrition is not a diet -- it IS care", 0),
    ("Food Forest: Fruit, berries, nuts, mushrooms -- picked fresh", 0),
    ("Raised Beds: Greens, peppers, squash -- reducing inflammation", 0),
    ("Hydroponic Lab: Year-round microgreens, zero pesticides", 0),
    ("Herb Garden: Lavender, chamomile, mint -- medicine from soil", 0),
    ("Eggs & Goat Milk: Protein from the estate animals", 0),
    ("Community Table: Elders, trainees, staff eat TOGETHER", 0),
    ("No one eats alone. Everyone eats what they grew.", 1),
])
print("8: Nutrition")

content("Spiritual Healing -- The Golden Years", [
    ("This is NOT a place to wait. This is a place to LIVE.", 0),
    ("Morning Ritual: Gather eggs, walk lavender, breathe", 0),
    ("Cedar Sauna: Soul healing, quiet time, storytelling", 0),
    ("Shared Meals: Food they grew together, around one table", 0),
    ("Evening Circles: Porch music, oral history, legacy", 0),
    ("Elders teach youth what schools can't -- wisdom, patience", 0),
    ("Peaceful Transition: Surrounded by love, on beautiful land", 0),
    ("Not in a cold institution -- at peace, with community", 1),
])
print("9: Spiritual Healing")

content("Community Service -- Serving Brunswick County", [
    ("This is for the lower-income community around them", 0),
    ("Brunswick County: high poverty, limited healthcare, aging", 1),
    ("Medicaid-Funded Beds -- No one turned away for income", 0),
    ("Community Meals -- Open farm-to-table for neighborhood seniors", 0),
    ("Free Wellness Days -- Cedar sauna, garden walks, hemp cream", 0),
    ("Youth Employment -- Local teens trained AND PAID", 0),
    ("Food Distribution -- Surplus produce to food-insecure families", 0),
    ("Every grant dollar serves the people who need it most", 1),
])
print("10: Community Service")

content("The Killer Grant Narrative", [
    ("Use this in EVERY application:", 0),
    ("", 0),
    ("\"We are building an integrated agrotherapy campus", 0),
    ("where elderly residents participate in lavender farming,", 1),
    ("hemp wellness production, and plant-based nutrition", 1),
    ("as evidence-based therapeutic programming,", 1),
    ("while training the next generation of rural healthcare", 1),
    ("workers and sustainable farmers --", 1),
    ("on historically significant Black-owned land", 1),
    ("in an underserved Virginia tobacco community.\"", 1),
])
print("6: Grant Narrative")

# =====================================================================
# PART 1B: THE COMPLETE SYSTEM
# =====================================================================
transition("The Circular System\nNothing Wasted, Everything Regenerated")
print("7: System transition")

content("The 8 Core Elements", [
    ("People -- Community, labor, education, mentorship", 0),
    ("Animals -- Goats, chickens, geese, alpaca on controlled routes", 0),
    ("Plants -- Lavender, 30-acre pine forest, hemp", 0),
    ("Soil -- Composting, mycelium, nutrient cycling", 0),
    ("Water -- Streams, wetlands, irrigation connecting all systems", 0),
    ("Electricity -- Off-grid solar, wind, hemp biomass fuel", 0),
    ("Internet -- E-commerce, training platform, crop monitoring", 0),
    ("Spirit -- Purpose, healing, legacy, generational vision", 0),
])
print("8: 8 Core Elements")

content("Program Elements -- The Puzzle Pieces", [
    ("The Homestead -- Residential HQ, command center of 66 acres", 0),
    ("Event Center -- Weddings, workshops, farm dinners, youth programs", 0),
    ("Animal Systems -- Goats, chickens, geese, alpaca managing the land", 0),
    ("Food Forest -- Permaculture feeding animals and people", 0),
    ("Hemp Innovation Hub -- 25,000+ products, NC A&T research", 0),
    ("The original \"Leaf\" in ReThinking The Leaf", 1),
])
print("9: Puzzle Pieces")

content("Hemp: 25,000+ Products -- Lamont's Vision", [
    ("Textiles & Clothing -- Denim, shoes, bags, upholstery", 0),
    ("Construction -- Hempcrete, insulation, fiberboard", 0),
    ("Paper & Printing -- Stronger than wood pulp", 0),
    ("Food & Nutrition -- Protein, oils, dairy alternatives", 0),
    ("Health & Body Care -- Soaps, lotions, lip balms", 0),
    ("Fuel & Industrial -- Biodiesel, ethanol from biomass", 0),
    ("Animal Products -- Bedding, pet treats, supplements", 0),
    ("Bioplastics -- Biodegradable plastics, car panels", 0),
])
print("10: Hemp Products")

content("11 Estate Zones", [
    ("Lavender Fields & Floralpy -- Sensory therapy, essential oils", 0),
    ("Agrotherapy Cabin Resort -- A-frame retreats, fire pits", 0),
    ("Agro-Hound Dog Park -- Support animals welcome", 0),
    ("Wetland Education Zone -- Boardwalks, outdoor classrooms", 0),
    ("Trail Network -- 30 acres of pine plantation trails", 0),
    ("Farm Store & Community Market -- Local produce & goods", 0),
    ("Hemp Fields & Phytoremediation -- Soil healing + fiber", 0),
    ("Hydroponic Lab, Butterfly Pavilion, Food Forest, Strolling Gardens", 0),
])
print("11: Estate Zones")

content("66 Acres -- Spatial Breakdown", [
    ("Pine Forest -- 30 acres (45%)", 0),
    ("Lavender Fields -- 8 acres (12%)", 0),
    ("Hemp Fields -- 5 acres (7.5%)", 0),
    ("Food Forest -- 5 acres (7.5%)", 0),
    ("Trails & Wetlands -- 5 acres (7.5%)", 0),
    ("Animal Areas -- 4 acres (6%)", 0),
    ("Event Center & House -- 3 acres", 0),
    ("Cabins & Resort -- 3 acres / Infrastructure -- 3 acres", 0),
])
print("12: Spatial Breakdown")

content("8 Training Programs -- Farm Academy 2027", [
    ("1. Off-Grid Solar & Energy -- panels, batteries, inverters", 0),
    ("2. Water Systems & Harvesting -- cisterns, drip, greywater", 0),
    ("3. Soil Science & Composting -- testing, biochar, cover crops", 0),
    ("4. Plant Processing & Permaculture -- distillation, hemp fiber", 0),
    ("5. Animal Husbandry -- grazing, poultry, beekeeping, butchering", 0),
    ("6. Community Building & Wellness -- co-ops, farm-to-table biz", 0),
    ("7. Natural Building & Off-Grid -- cob, timber, A-frames", 0),
    ("8. Zero-Waste Farm Ops -- vermiculture, biodigester, closed-loop", 0),
])
print("13: Training Programs")

# =====================================================================
# PART 2: FARM GRANT STRATEGY
# =====================================================================
transition("Part 2: Farm Grant Strategy\n$925K+ Identified")
print("14: Farm Grants transition")

content("URGENT Priority Grants", [
    ("AFID Infrastructure Grant (Virginia AFID) -- Up to $50,000", 0),
    ("Agrotherapy Processing Pavilion & Farm Store", 1),
    ("VAPG Value-Added Producer (USDA AMS) -- Up to $75,000", 0),
    ("Lavender soaps, oils, farm store, e-commerce", 1),
    ("VA BIPOC Historic Preservation (Virginia DHR) -- Varies", 0),
    ("Historic site stabilization, Lawrenceville Historic District", 1),
])
print("15: URGENT Grants")

content("HIGH Priority Grants", [
    ("TRIAD -- Tobacco Region Incentive (VA TRRC) -- Varies", 0),
    ("Brunswick County is eligible tobacco region county", 1),
    ("EQIP Environmental Quality (USDA NRCS) -- Up to $450,000", 0),
    ("Wetland buffers, pollinator habitats, soil conservation", 1),
    ("Socially-Disadvantaged Groups (USDA FSA) -- Varies", 0),
    ("BIPOC landowner support, minority ownership structure", 1),
])
print("16: HIGH Grants")

content("MEDIUM & PLANNED Grants", [
    ("VSU Small Farm Outreach -- Free Technical Assistance", 0),
    ("Virginia BMP Tax Credit -- 25% of BMP Costs", 0),
    ("Beginning Farmer & Rancher (USDA NIFA) -- Up to $250,000", 0),
    ("Specialty Crop Block Grant (VDACS) -- Targeting 2027", 0),
    ("USDA REAP Rural Energy -- Solar array consideration", 0),
])
print("17: MEDIUM Grants")

# =====================================================================
# PART 3: ELDERLY CARE GRANTS
# =====================================================================
transition("Part 3: Elderly Care Facility\nGrant Opportunities")
print("18: Elder Care transition")

content("Top Priority -- Federal Programs", [
    ("USDA Community Facilities Grant -- Up to 75% of costs", 0),
    ("FY2026: ~$671M nationally, rolling applications", 1),
    ("Must be nonprofit (consider 501(c)(3) structure)", 1),
    ("Rural Health Transformation Program -- $50B over 5 years", 0),
    ("Brand new 2026-2030 program via CMS", 1),
    ("VA Tobacco Region Commission -- $1.35M+ recent awards", 0),
    ("Brunswick County is in service area", 1),
])
print("19: Federal Programs")

content("Virginia State & Medicaid Programs", [
    ("Virginia Auxiliary Grant -- Revenue stream for approved ALFs", 0),
    ("CCC Plus Waiver -- Recurring Medicaid revenue", 0),
    ("Must enroll as provider through DMAS", 1),
    ("DARS Programs -- Respite vouchers, independent living", 0),
    ("VA Aid & Attendance -- For eligible veteran residents", 0),
])
print("20: State Programs")

content("BIPOC-Specific Opportunities", [
    ("SBA 8(a) Business Development -- Strengthens federal apps", 0),
    ("New Markets Tax Credit -- 39%, Brunswick County eligible", 0),
    ("MBDA Business Centers -- Free consulting for loan packages", 0),
    ("CDFIs (The Innovate Fund) -- Health/wellness priority", 0),
    ("RWJF & AARP Foundations -- Health equity grants", 0),
])
print("21: BIPOC Opportunities")

content("Virginia Licensing -- ALFA Requirements", [
    ("Licensed by: VA Dept of Social Services (VDSS)", 0),
    ("Administrator License (ALFA) required", 0),
    ("Training: 320-640 hours (education dependent)", 1),
    ("Exam: National RC/AL examination (NAB)", 1),
    ("Fees: ~$595 total", 1),
    ("Facility: VDSS inspection, background checks, TB screening", 0),
    ("Apply via VELA Portal: licensing.dss.virginia.gov", 0),
])
print("22: Licensing")

# =====================================================================
# PART 4: CHRONIC GURU
# =====================================================================
transition("Part 4: Chronic Guru\nBusiness Model Analysis")
print("23: Chronic Guru transition")

content("Chronic Guru -- What Is It?", [
    ("Cannabis/hemp dispensary brand (NOT health/wellness)", 0),
    ("Founded 2016 by Patrick O'Brien (veteran)", 0),
    ("6 locations: 4 in Florida, 2 in North Carolina", 0),
    ("Exploits THCa/hemp loophole under 2018 Farm Bill", 0),
    ("3 revenue streams: retail lounges, e-commerce, franchising", 0),
])
print("24: Chronic Guru overview")

content("Critical Warning -- Legal Loophole Closing", [
    ("H.R. 5371 signed November 2025", 0),
    ("Shifts to \"total THC\" standard (Delta-9 + THCa)", 0),
    ("EFFECTIVE: November 12, 2026 -- 5 months away", 0),
    ("After this date: THCa products = marijuana", 0),
    ("Recommendation: DO NOT pursue this model", 0),
])
print("25: Legal Warning")

# =====================================================================
# RECOMMENDATIONS & CLOSE
# =====================================================================
content("Strategic Recommendations", [
    ("Register on SAM.gov immediately (takes weeks)", 0),
    ("Contact USDA VA Rural Development for CF Grant guidance", 0),
    ("Contact VA Tobacco Region Commission", 0),
    ("Consider 501(c)(3) for elder care component", 0),
    ("Begin ALFA administrator license pathway", 0),
    ("Position as \"Agrotherapy + Elder Care\" in all applications", 0),
])
print("26: Recommendations")

content("Revenue Flywheel -- Self-Sustaining Model", [
    ("LAVENDER FARM produces PRODUCTS", 0),
    ("Sells to: Farm Store, E-commerce, Elder Facility", 1),
    ("ELDER CARE FACILITY funded by:", 0),
    ("Medicaid 40% / VA Auxiliary 25% / Private Pay 20%", 1),
    ("Veterans 10% / Grants 5%", 1),
    ("TRAINING PROGRAMS produce workforce for BOTH", 0),
    ("Farm feeds facility. Facility funds farm.", 0),
    ("Training programs staff both.", 0),
])
print("27: Revenue Flywheel")

transition("Questions & Next Steps")
print("28: Closing")

# =====================================================================
# SAVE
# =====================================================================
prs.save(OUTPUT)
print(f"\nSaved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
